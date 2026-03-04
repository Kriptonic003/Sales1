"""
rag_service.py — One-shot RAG pipeline (multi-format, TF-IDF retrieval)
Supported: PDF, DOCX, XLSX/XLS, CSV, PPTX, TXT, MD
LLM priority: Groq (primary) → Gemini raw HTTP (fallback)

Uses sklearn TF-IDF + cosine similarity for retrieval.
No HuggingFace embedding model — avoids process crash on Windows.
"""

import csv
import io
import logging
import os
from typing import Tuple

import requests as _http

logger = logging.getLogger(__name__)

# ── Gemini REST endpoint ──────────────────────────────────────────────────────
_GEMINI_URL = (
    "https://generativelanguage.googleapis.com"
    "/v1beta/models/gemini-1.5-flash:generateContent"
)

# ── Supported formats ─────────────────────────────────────────────────────────
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".csv", ".pptx", ".txt", ".md"
}


# ── LLM helpers ───────────────────────────────────────────────────────────────
def _call_groq(prompt: str, api_key: str) -> str:
    from langchain_groq import ChatGroq
    from langchain_core.messages import HumanMessage
    llm = ChatGroq(model="llama-3.1-8b-instant", groq_api_key=api_key, temperature=0.3)
    result = llm.invoke([HumanMessage(content=prompt)])
    return result.content


def _call_gemini(prompt: str, api_key: str) -> str:
    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {"temperature": 0.3},
    }
    resp = _http.post(
        _GEMINI_URL, params={"key": api_key}, json=payload, timeout=60
    )
    resp.raise_for_status()
    data = resp.json()
    return (
        data.get("candidates", [{}])[0]
        .get("content", {})
        .get("parts", [{}])[0]
        .get("text", "")
    ) or "No answer generated."


def _call_llm(prompt: str) -> Tuple[str, str]:
    """Try Groq first, fall back to Gemini. Returns (answer, llm_name)."""
    groq_key = os.getenv("GROQ_API_KEY", "").strip()
    gemini_key = os.getenv("GEMINI_API_KEY", "").strip()

    if groq_key:
        try:
            return _call_groq(prompt, groq_key), "groq"
        except Exception as e:
            logger.warning(f"Groq failed: {e}")

    if gemini_key:
        try:
            return _call_gemini(prompt, gemini_key), "gemini"
        except Exception as e:
            logger.warning(f"Gemini failed: {e}")

    raise RuntimeError(
        "No LLM available. Add GROQ_API_KEY to .env (free at https://console.groq.com)."
    )


# ── Multi-format file parser ──────────────────────────────────────────────────
def parse_file(file_bytes: bytes, filename: str) -> str:
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            "Supported: PDF, DOCX, XLSX, CSV, PPTX, TXT, MD"
        )

    text = ""

    if ext == ".pdf":
        import pdfplumber
        parts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    parts.append(t.strip())
        text = "\n\n".join(parts)

    elif ext in (".docx", ".doc"):
        import docx
        doc = docx.Document(io.BytesIO(file_bytes))
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    parts.append(row_text)
        text = "\n\n".join(parts)

    elif ext in (".xlsx", ".xls"):
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                row_str = " | ".join(
                    str(c) for c in row if c is not None and str(c).strip()
                )
                if row_str:
                    parts.append(row_str)
        text = "\n".join(parts)

    elif ext == ".csv":
        decoded = file_bytes.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(decoded))
        rows = [
            " | ".join(cell.strip() for cell in row)
            for row in reader if any(c.strip() for c in row)
        ]
        text = "\n".join(rows)

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        text = "\n\n".join(parts)

    elif ext in (".txt", ".md"):
        text = file_bytes.decode("utf-8", errors="replace")

    if not text.strip():
        raise ValueError(
            f"Could not extract any text from '{filename}'. "
            "The file may be empty, image-based, or corrupted."
        )
    return text.strip()


# ── TF-IDF retrieval (no ML model — avoids Windows process crash) ─────────────
def _retrieve_top_k(chunks: list, question: str, k: int = 5) -> list:
    """Use TF-IDF cosine similarity to find the top-k relevant chunks."""
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity

    corpus = chunks + [question]
    vectorizer = TfidfVectorizer(stop_words="english", max_features=5000)
    tfidf = vectorizer.fit_transform(corpus)

    query_vec = tfidf[-1]          # last entry = question
    chunk_vecs = tfidf[:-1]        # everything else = chunks
    scores = cosine_similarity(query_vec, chunk_vecs).flatten()

    top_idx = scores.argsort()[::-1][:k]
    return [chunks[i] for i in top_idx if scores[i] > 0] or chunks[:k]


def _split_text(text: str, max_chunks: int = 150) -> list:
    """Split text into chunks using simple line-based strategy."""
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    # Group lines into chunks of ~10 lines each
    chunk_size = 10
    chunks = []
    for i in range(0, len(lines), chunk_size):
        chunk = "\n".join(lines[i:i + chunk_size])
        if chunk:
            chunks.append(chunk)
    if not chunks:
        raise ValueError("Text could not be split into chunks.")
    if len(chunks) > max_chunks:
        logger.info(f"Capping {len(chunks)} chunks → {max_chunks}")
        chunks = chunks[:max_chunks]
    return chunks


# ── Main RAG entry point ──────────────────────────────────────────────────────
def run_rag(file_bytes: bytes, filename: str, question: str) -> dict:
    """Full RAG pipeline. Returns { answer, llm_used, source_chunks }."""

    # 1. Parse
    full_text = parse_file(file_bytes, filename)

    # 2. Chunk
    chunks = _split_text(full_text)

    # 3. Retrieve top-5 by TF-IDF similarity
    source_chunks = _retrieve_top_k(chunks, question, k=5)

    # 4. Build prompt
    context = "\n\n---\n\n".join(source_chunks)
    prompt = (
        "You are a sales analytics expert. Use ONLY the context below to answer "
        "the question. If the answer is not in the context, say 'I could not find "
        "that information in the uploaded document.'\n\n"
        f"Context:\n{context}\n\n"
        f"Question: {question}\n\n"
        "Answer (be concise and business-friendly):"
    )

    # 5. Call LLM
    answer, llm_name = _call_llm(prompt)

    return {
        "answer": answer,
        "llm_used": llm_name,
        "source_chunks": [c[:300] for c in source_chunks],
    }
